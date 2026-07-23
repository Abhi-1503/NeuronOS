import io

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


class UnsupportedFileTypeError(Exception):
    pass


def extract_text(*, file_type: str, content: bytes) -> str:
    """Real text extraction, not a stub — genuinely reads the file's content so keyword
    search (Roadmap Phase 1's actual Knowledge scope) has something real to search over,
    independent of whether an LLM summarization key is configured."""
    if file_type == "pdf":
        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if file_type == "docx":
        doc = DocxDocument(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)
    if file_type == "pptx":
        prs = Presentation(io.BytesIO(content))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    if file_type == "xlsx":
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                lines.append(" ".join(str(cell) for cell in row if cell is not None))
        return "\n".join(lines)
    if file_type in ("email", "other"):
        # Best-effort plain-text decode — email/"other" aren't a fixed binary format,
        # so there's no dedicated parser to reach for the way there is for the above.
        try:
            return content.decode("utf-8", errors="ignore")
        except Exception:
            return ""
    raise UnsupportedFileTypeError(file_type)


def chunk_text(text: str, *, chunk_size_chars: int = 2000) -> list[str]:
    """~500 tokens per chunk (Blueprint §14.2's RAG pipeline) — approximated as ~4
    chars/token, a well-known rough English-text ratio, not an exact tokenizer count;
    exact-enough for chunk sizing, not claimed as precise."""
    text = text.strip()
    if not text:
        return []
    return [text[i : i + chunk_size_chars] for i in range(0, len(text), chunk_size_chars)]
